"""
Presentación Ejecutiva - Análisis de Modelo Predictivo y Proveedores
Solventa - Prueba Técnica Data Scientist Jr.
"""

import warnings

warnings.filterwarnings("ignore")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json

# Load metrics
with open("output/model_metrics.json", "r") as f:
    model_metrics = json.load(f)
with open("output/competitor_metrics.json", "r") as f:
    competitor_metrics = json.load(f)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(0x2C, 0x3E, 0x50)
BLUE = RGBColor(0x34, 0x98, 0xDB)
LIGHT_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)


def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=DARK_TEXT,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    font_size=14,
    color=DARK_TEXT,
    spacing=Pt(8),
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_header_bar(slide, title_text):
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), NAVY)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.25),
        Inches(11),
        Inches(0.7),
        title_text,
        font_size=28,
        color=WHITE,
        bold=True,
    )
    # Thin accent line
    add_shape(slide, Inches(0), Inches(1.2), prs.slide_width, Inches(0.05), BLUE)


def add_slide_number(slide, num, total=10):
    add_text_box(
        slide,
        Inches(12.2),
        Inches(7.0),
        Inches(0.8),
        Inches(0.4),
        f"{num}/{total}",
        font_size=10,
        color=GRAY,
        alignment=PP_ALIGN.RIGHT,
    )


# ====================
# SLIDE 1: COVER
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape(slide, Inches(0), Inches(3.5), prs.slide_width, Inches(0.06), BLUE)

add_text_box(
    slide,
    Inches(1.5),
    Inches(1.5),
    Inches(10),
    Inches(1.2),
    "Solventa",
    font_size=52,
    color=WHITE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

add_text_box(
    slide,
    Inches(1.5),
    Inches(3.8),
    Inches(10),
    Inches(1.5),
    "Análisis de Modelo Predictivo\ny Evaluación de Proveedores",
    font_size=32,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.CENTER,
)

add_text_box(
    slide,
    Inches(1.5),
    Inches(5.5),
    Inches(10),
    Inches(0.8),
    "Prueba Técnica - Data Scientist Jr.  |  Mayo 2026",
    font_size=16,
    color=RGBColor(0xB0, 0xBE, 0xC5),
    alignment=PP_ALIGN.CENTER,
)

# ====================
# SLIDE 2: AGENDA
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Agenda")
add_slide_number(slide, 1)

agenda_items = [
    ("01", "Modelo Predictivo Interno", "Construcción, evaluación y punto de corte"),
    ("02", "Evaluación de Proveedores", "Análisis comparativo AB vs XY"),
    ("03", "Conclusiones y Recomendaciones", "Decisiones estratégicas para el negocio"),
]

y_start = Inches(2.0)
for num, title, desc in agenda_items:
    add_shape(slide, Inches(1.5), y_start, Inches(0.8), Inches(0.8), BLUE)
    add_text_box(
        slide,
        Inches(1.5),
        y_start + Inches(0.15),
        Inches(0.8),
        Inches(0.5),
        num,
        font_size=24,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(2.6),
        y_start + Inches(0.05),
        Inches(8),
        Inches(0.4),
        title,
        font_size=22,
        color=NAVY,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(2.6),
        y_start + Inches(0.45),
        Inches(8),
        Inches(0.3),
        desc,
        font_size=14,
        color=GRAY,
    )
    y_start += Inches(1.3)

# ====================
# SLIDE 3: CONTEXTO
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Contexto del Análisis")
add_slide_number(slide, 2)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.6),
    Inches(5.5),
    Inches(5),
    "La compañía evalúa el lanzamiento de un nuevo producto dirigido a un "
    "segmento de clientes de alto riesgo.",
    font_size=20,
    color=DARK_TEXT,
)

add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(2.6),
    Inches(5.5),
    Inches(4),
    [
        "• 4,315 clientes potenciales analizados",
        "• 22 variables disponibles (demográficas, financieras)",
        "• Tasa de mora 30+ días: 24.4%",
        "• Tasa de mora 60+ días: 11.7%",
        "• 2 proveedores evaluados (AB y XY)",
    ],
    font_size=16,
)

# Key metrics boxes
metrics = [
    ("4,315", "Clientes\nAnalizados"),
    ("24.4%", "Tasa Mora\n30+ días"),
    ("114K", "Registros\nProveedores"),
]

for i, (val, label) in enumerate(metrics):
    x = Inches(7.5 + i * 2.0)
    add_shape(slide, x, Inches(2.5), Inches(1.7), Inches(2.2), LIGHT_BLUE)
    add_text_box(
        slide,
        x,
        Inches(2.7),
        Inches(1.7),
        Inches(0.8),
        val,
        font_size=36,
        color=BLUE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        x,
        Inches(3.5),
        Inches(1.7),
        Inches(1),
        label,
        font_size=13,
        color=NAVY,
        alignment=PP_ALIGN.CENTER,
    )

# ====================
# SLIDE 4: MODELO PREDICTIVO - RESULTADOS
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "1. Modelo Predictivo - Resultados")
add_slide_number(slide, 3)

add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(1.5),
    Inches(5.5),
    Inches(3),
    [
        "• Variable objetivo: Mora30 (detección temprana)",
        "• Modelo seleccionado: Logistic Regression",
        "• ROC-AUC: 0.61 (poder predictivo moderado)",
        "• Punto de corte: 0.49",
        "• Se captura 59% de clientes morosos",
        "• Tasa de aprobación: 53%",
    ],
    font_size=16,
)

# Key metrics
box_data = [
    ("ROC-AUC", f"{model_metrics['roc_auc']:.2f}", "Capacidad\ndiscriminante"),
    ("Recall", f"{model_metrics['recall'] * 100:.0f}%", "Morosos\ndetectados"),
    (
        "Aprobación",
        f"{model_metrics['approval_rate'] * 100:.0f}%",
        "Clientes\naprobados",
    ),
]

for i, (metric, val, label) in enumerate(box_data):
    x = Inches(7.5 + i * 1.9)
    add_shape(slide, x, Inches(2.0), Inches(1.6), Inches(1.8), LIGHT_BLUE)
    add_text_box(
        slide,
        x,
        Inches(2.15),
        Inches(1.6),
        Inches(0.4),
        metric,
        font_size=14,
        color=NAVY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        x,
        Inches(2.5),
        Inches(1.6),
        Inches(0.6),
        val,
        font_size=30,
        color=BLUE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        x,
        Inches(3.1),
        Inches(1.6),
        Inches(0.6),
        label,
        font_size=11,
        color=GRAY,
        alignment=PP_ALIGN.CENTER,
    )

# Add model comparison chart image
img_left = Inches(0.5)
img_top = Inches(4.2)
try:
    slide.shapes.add_picture(
        "output/figures/04_roc_pr_curves.png",
        img_left,
        img_top,
        width=Inches(6),
        height=Inches(2.8),
    )
except:
    pass

try:
    slide.shapes.add_picture(
        "output/figures/07_confusion_matrix.png",
        Inches(7.0),
        Inches(4.3),
        width=Inches(2.8),
        height=Inches(2.3),
    )
except:
    pass

# ====================
# SLIDE 5: PUNTO DE CORTE
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "1.1 Punto de Corte Recomendado")
add_slide_number(slide, 4)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.5),
    Inches(11),
    Inches(0.6),
    "Para un producto de alto riesgo, el punto de corte balancea captura de morosos con viabilidad comercial.",
    font_size=16,
    color=GRAY,
)

# Cut-off recommendation box
add_shape(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.0), NAVY)
add_text_box(
    slide,
    Inches(1.0),
    Inches(2.4),
    Inches(5),
    Inches(0.4),
    "PUNTO DE CORTE RECOMENDADO",
    font_size=14,
    color=BLUE,
    bold=True,
)
add_text_box(
    slide,
    Inches(1.0),
    Inches(2.8),
    Inches(5),
    Inches(0.7),
    "0.49",
    font_size=48,
    color=WHITE,
    bold=True,
)
add_text_box(
    slide,
    Inches(1.0),
    Inches(3.6),
    Inches(5),
    Inches(0.5),
    "Cliente con probabilidad ≥ 49% de mora → RECHAZAR",
    font_size=14,
    color=RGBColor(0xB0, 0xBE, 0xC5),
)

# Impact metrics
impact_data = [
    ("Precision", "30.5%", "De rechazados,\n30% son morosos reales"),
    ("Recall", "58.6%", "Se detecta 59% de\nlos morosos"),
    ("Approval", "53.2%", "Poco más de la mitad\nde clientes aprobados"),
]

for i, (metric, val, desc) in enumerate(impact_data):
    y = Inches(2.3)
    x = Inches(6.8 + i * 2.2)
    color = GREEN if metric == "Recall" else (RED if metric == "Approval" else BLUE)
    add_shape(slide, x, y, Inches(2.0), Inches(2.0), LIGHT_BLUE)
    add_text_box(
        slide,
        x,
        y + Inches(0.15),
        Inches(2),
        Inches(0.3),
        metric,
        font_size=14,
        color=NAVY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        x,
        y + Inches(0.5),
        Inches(2),
        Inches(0.5),
        val,
        font_size=28,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        x,
        y + Inches(1.1),
        Inches(2),
        Inches(0.8),
        desc,
        font_size=11,
        color=GRAY,
        alignment=PP_ALIGN.CENTER,
    )

try:
    slide.shapes.add_picture(
        "output/figures/06_cutoff_analysis.png",
        Inches(0.5),
        Inches(4.5),
        width=Inches(12),
        height=Inches(2.5),
    )
except:
    pass

# ====================
# SLIDE 6: PROVEEDORES - COMPARATIVA
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "2. Evaluación de Proveedores - Comparativa")
add_slide_number(slide, 5)

add_text_box(
    slide,
    Inches(0.8),
    Inches(1.5),
    Inches(11),
    Inches(0.6),
    "Análisis de 114,109 registros con puntajes de dos proveedores y variable de default (60.4% tasa)",
    font_size=16,
    color=GRAY,
)

# Comparison table
comp_data = [
    ["Métrica", "Proveedor AB", "Proveedor XY", "Benchmark"],
    [
        "ROC-AUC",
        f"{competitor_metrics['auc_ab']:.4f}",
        f"{competitor_metrics['auc_xy']:.4f}",
        "> 0.65 (aceptable)",
    ],
    [
        "KS Statistic",
        f"{competitor_metrics['ks_ab']:.4f}",
        f"{competitor_metrics['ks_xy']:.4f}",
        "> 0.20 (bueno)",
    ],
    [
        "Gini",
        f"{competitor_metrics['gini_ab']:.4f}",
        f"{competitor_metrics['gini_xy']:.4f}",
        "> 0.30",
    ],
    [
        "PSI",
        f"{competitor_metrics['psi_ab']:.4f}",
        f"{competitor_metrics['psi_xy']:.4f}",
        "< 0.10 (estable)",
    ],
    ["Cobertura", "100%", "88.8%", "100%"],
]

table = slide.shapes.add_table(
    len(comp_data), 4, Inches(1.5), Inches(2.3), Inches(10.3), Inches(2.5)
).table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.5)
table.columns[3].width = Inches(3.3)

for i, row in enumerate(comp_data):
    for j, cell_text in enumerate(row):
        cell = table.cell(i, j)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.name = "Calibri"
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            p.font.color.rgb = DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if i % 2 == 0 else WHITE
        p.alignment = PP_ALIGN.CENTER

try:
    slide.shapes.add_picture(
        "output/figures/09_roc_competencia.png",
        Inches(1.0),
        Inches(4.9),
        width=Inches(5.5),
        height=Inches(2.2),
    )
except:
    pass

try:
    slide.shapes.add_picture(
        "output/figures/10_deciles_competencia.png",
        Inches(6.8),
        Inches(4.9),
        width=Inches(5.5),
        height=Inches(2.2),
    )
except:
    pass

# ====================
# SLIDE 7: PROVEEDORES - HALLAZGOS
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "2.1 Hallazgos Críticos - Proveedores")
add_slide_number(slide, 6)

# Big warning
add_shape(
    slide,
    Inches(1.0),
    Inches(1.6),
    Inches(11.3),
    Inches(1.0),
    RGBColor(0xFD, 0xED, 0xEC),
)
add_text_box(
    slide,
    Inches(1.3),
    Inches(1.75),
    Inches(10.5),
    Inches(0.7),
    "⚠ AMBOS MODELOS TIENEN PODER DISCRIMINANTE EQUIVALENTE AL AZAR (AUC ≈ 0.50)",
    font_size=18,
    color=RED,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(3.0),
    Inches(5.5),
    Inches(4),
    [
        "• La tasa de default es ~60% en todos los deciles",
        "• No hay monotonicidad en la clasificación",
        "• Los puntajes no separan buenos de malos",
        "• KS < 0.01 indica separación nula",
        "• Gini ≈ 0 confirma falta de poder predictivo",
    ],
    font_size=16,
)

# Visual evidence
try:
    slide.shapes.add_picture(
        "output/figures/13_separation_analysis.png",
        Inches(6.8),
        Inches(3.0),
        width=Inches(5.5),
        height=Inches(2.5),
    )
except:
    pass

try:
    slide.shapes.add_picture(
        "output/figures/11_default_rate_deciles.png",
        Inches(6.8),
        Inches(5.5),
        width=Inches(5.5),
        height=Inches(1.5),
    )
except:
    pass

# ====================
# SLIDE 8: RECOMENDACIÓN PROVEEDORES
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "2.2 Recomendación sobre Proveedores")
add_slide_number(slide, 7)

# Main recommendation
add_shape(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.5), NAVY)
add_text_box(
    slide,
    Inches(1.5),
    Inches(1.9),
    Inches(10.5),
    Inches(0.5),
    "RECOMENDACIÓN",
    font_size=16,
    color=BLUE,
    bold=True,
)
add_text_box(
    slide,
    Inches(1.5),
    Inches(2.3),
    Inches(10.5),
    Inches(0.8),
    "No contratar a ninguno de los dos proveedores",
    font_size=30,
    color=WHITE,
    bold=True,
)

# Reasons
add_text_box(
    slide,
    Inches(0.8),
    Inches(3.6),
    Inches(5.5),
    Inches(0.4),
    "Razones:",
    font_size=18,
    color=NAVY,
    bold=True,
)

add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(4.1),
    Inches(5.5),
    Inches(3),
    [
        "• Sin capacidad discriminante (AUC ≈ 0.50)",
        "• Costo sin retorno: equivalente a seleccionar al azar",
        "• Proveedor XY: 11.2% de datos sin puntaje",
        "• Rangos de scores muy estrechos (AB: std=0.65)",
        "• No generarían valor vs. selección aleatoria",
    ],
    font_size=15,
)

# Alternatives
add_shape(slide, Inches(7.0), Inches(3.6), Inches(5.3), Inches(3.5), LIGHT_BLUE)
add_text_box(
    slide,
    Inches(7.3),
    Inches(3.7),
    Inches(4.8),
    Inches(0.4),
    "Alternativas Recomendadas",
    font_size=18,
    color=NAVY,
    bold=True,
)

add_bullet_slide(
    slide,
    Inches(7.3),
    Inches(4.2),
    Inches(4.8),
    Inches(2.8),
    [
        "1. Desarrollar modelo interno con mejores datos",
        "2. Incorporar buró de crédito (historial, score)",
        "3. Usar datos transaccionales y de comportamiento",
        "4. Explorar proveedores alternativos",
        "5. Implementar scoring híbrido (interno + externo)",
    ],
    font_size=14,
    color=DARK_TEXT,
)

# ====================
# SLIDE 9: CONCLUSIONES
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "3. Conclusiones y Plan de Acción")
add_slide_number(slide, 8)

# Conclusions left
add_text_box(
    slide,
    Inches(0.8),
    Inches(1.5),
    Inches(5.5),
    Inches(0.4),
    "Conclusiones",
    font_size=20,
    color=NAVY,
    bold=True,
)

add_bullet_slide(
    slide,
    Inches(0.8),
    Inches(2.1),
    Inches(5.5),
    Inches(5),
    [
        "✓ Modelo interno con AUC=0.61: punto de partida útil",
        "✓ Punto de corte 0.49: balancea riesgo y negocio",
        "✗ Proveedores AB y XY: sin poder predictivo",
        "✗ Datos actuales: insuficientes para alto rendimiento",
        "⚠ Riesgo residual significativo en aprobaciones",
    ],
    font_size=15,
    color=DARK_TEXT,
    spacing=Pt(12),
)

# Action plan right
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), NAVY)
add_text_box(
    slide,
    Inches(7.1),
    Inches(1.7),
    Inches(5.3),
    Inches(0.4),
    "PLAN DE ACCIÓN",
    font_size=18,
    color=BLUE,
    bold=True,
)

actions = [
    ("Corto plazo", "Implementar modelo actual con monitoreo estricto"),
    ("Mediano plazo", "Integrar datos de buró y transaccionales"),
    ("Largo plazo", "Desarrollar modelo interno robusto y automatizado"),
]

for i, (term, action) in enumerate(actions):
    y = Inches(2.3 + i * 1.5)
    add_shape(
        slide, Inches(7.1), y, Inches(5.3), Inches(1.3), RGBColor(0x34, 0x49, 0x5E)
    )
    add_text_box(
        slide,
        Inches(7.3),
        y + Inches(0.1),
        Inches(2),
        Inches(0.3),
        term,
        font_size=14,
        color=BLUE,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(7.3),
        y + Inches(0.45),
        Inches(4.9),
        Inches(0.7),
        action,
        font_size=13,
        color=WHITE,
    )

# ====================
# SLIDE 10: CIERRE
# ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape(slide, Inches(0), Inches(3.5), prs.slide_width, Inches(0.06), BLUE)

add_text_box(
    slide,
    Inches(1.5),
    Inches(2.0),
    Inches(10),
    Inches(1),
    "Gracias",
    font_size=52,
    color=WHITE,
    bold=True,
    alignment=PP_ALIGN.CENTER,
)

add_text_box(
    slide,
    Inches(1.5),
    Inches(3.8),
    Inches(10),
    Inches(1),
    "Próximos pasos: Enriquecer datos → Mejorar modelo → Reducir riesgo",
    font_size=18,
    color=RGBColor(0xB0, 0xBE, 0xC5),
    alignment=PP_ALIGN.CENTER,
)

add_text_box(
    slide,
    Inches(1.5),
    Inches(5.0),
    Inches(10),
    Inches(1),
    "Solventa  |  Bogotá, Colombia  |  www.solventa.co",
    font_size=14,
    color=GRAY,
    alignment=PP_ALIGN.CENTER,
)

# Save
prs.save("output/presentacion_ejecutiva.pptx")
print("Presentación guardada: output/presentacion_ejecutiva.pptx")
